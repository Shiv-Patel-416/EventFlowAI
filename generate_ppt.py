import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]
picture_slide_layout = prs.slide_layouts[5] # Title only

# Helper to format titles
def format_title(title_shape, text):
    title_shape.text = text
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.name = 'Arial'
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(6, 182, 212) # Cyan color from the branding

# 1. Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "EventFlow AI"
subtitle.text = "Traffic Intelligence & Orchestration Platform\nProactive, AI-driven urban mobility management"
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(6, 182, 212)

# 2. The Problem
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
format_title(title_shape, "The Motive: The Problem")
tf = body_shape.text_frame
tf.text = "Current traffic management solutions are REACTIVE"
p = tf.add_paragraph()
p.text = "Political rallies, festivals, sports events, construction activities, and sudden VIP movements create localized traffic breakdowns."
p.level = 1
p = tf.add_paragraph()
p.text = "Police only respond to congestion AFTER it happens."
p.level = 1

# 3. The Solution
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
format_title(title_shape, "The Solution: EventFlow AI")
tf = body_shape.text_frame
tf.text = "A PROACTIVE full-stack intelligence platform"
p = tf.add_paragraph()
p.text = "Fuses Machine Learning with Operations Research"
p.level = 1
p = tf.add_paragraph()
p.text = "Predicts traffic disruptions BEFORE they peak"
p.level = 1
p = tf.add_paragraph()
p.text = "Automatically calculates congestion severity, duration, and the exact resource deployment plan (Police & Barricades)."
p.level = 1

# Helper for screenshot slides
def add_screenshot_slide(title_text, img_path, description=""):
    slide = prs.slides.add_slide(picture_slide_layout)
    format_title(slide.shapes.title, title_text)
    
    # Add image
    if os.path.exists(img_path):
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
    
    if description:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(100, 100, 100)

# 4. Overview
add_screenshot_slide("Command Center / Overview", "presentation/dashboard.png", "Real-time tracking of active incidents with high-level KPIs for city managers.")

# 5. Live Map
add_screenshot_slide("Geospatial Intel & Route Diversion", "presentation/map.png", "A*/Dijkstra algorithms calculate fastest alternative routes around incidents to prevent cascade jams.")

# 6. AI Predictor
add_screenshot_slide("AI Predictor (Impact Simulation)", "presentation/predictions.png", "XGBoost v3.0 calculates Severity Score, Est. Duration, and Resource Plan based on event parameters.")

# 7. Analytics
add_screenshot_slide("Jurisdiction Leaderboard & Analytics", "presentation/analytics.png", "Gamified tracking for Police Stations based on dynamic Efficiency Score.")

# 8. How it works
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
format_title(title_shape, "Under The Hood")
tf = body_shape.text_frame
tf.text = "Data-Driven & Dynamic"
p = tf.add_paragraph()
p.text = "Nothing is hard-coded. XGBoost v3.0 ML models trained on 8,173 real-world traffic records."
p.level = 1
p = tf.add_paragraph()
p.text = "Dynamic Efficiency Score = (Actual Resolution Time) / (AI Baseline Predicted Time)"
p.level = 1
p = tf.add_paragraph()
p.text = "Location Variance: Predictions change purely based on geographical features (e.g. MG Road vs small side street)."
p.level = 1

# 9. Team
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
format_title(title_shape, "Team: techy Boys")
tf = body_shape.text_frame
tf.text = "Shiv Patel (Leader)"
p = tf.add_paragraph()
p.text = "Saahil Saitwal"
p = tf.add_paragraph()
p.text = "Ronak Rathod"

# 10. Thank You
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You"
subtitle.text = "EventFlow AI - Redefining Urban Mobility"
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(6, 182, 212)

# Save
prs.save("presentation/EventFlowAI_Pitch.pptx")
print("Presentation generated successfully!")
